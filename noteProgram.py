# Author : Kevin De La Torre
# Purpose: Take a ".notes" file and convert it to a powerpoint file using "python-pptx"

from sys import argv
import os
from pptx import Presentation

PRES_FOLDER_NAME = "presentations"
SLD_TITLE_PRESENTATION = 0
SLD_TITLE_CONTENT = 1
SLD_SECTION_HEADER = 2
SLD_TITLE_ONLY = 5
SLD_BLANK = 6

def main():
    if ( len( argv ) > 1 ):
        folderPath = os.getcwd() + '/' + PRES_FOLDER_NAME
        if not os.path.exists( folderPath ):
                os.makedirs( folderPath )

        for i in range( 1, len( argv ) ): # Go through all files
            # Setup for presentation
            file = open( argv[ i ], 'r' ) # Open file in read-only mode
            prs = Presentation()
            slide = prs.slides.add_slide( prs.slide_layouts[ SLD_TITLE_PRESENTATION ] ) # Create presentation title slide
            title = slide.shapes.title
            title.text = argv[ i ].split( '/', 1 )[ 1 ].rsplit( '_', 1 )[ 0 ].replace( '_', ' ' )

            parser( prs, file )

            fileFolderPath = folderPath + '/' + argv[ i ].split( '/', 1 )[ 0 ] 

            # Creates a subfolder to store finished .pptxs
            if not os.path.exists( fileFolderPath ):
                os.makedirs( fileFolderPath )

            # Gets rid of folder name at left then replace file extension and places in folder
            if len( argv[ i ].split( '/' ) ) == 1:
                fileLocation = fileForderPath + '/' + argv[ i ].rsplit( '.', 1 )[ 0 ] + ".pptx"
            else:
                fileLocation = fileFolderPath + '/' + argv[ i ].split( '/', 1 )[ 1 ].rsplit( '.', 1 )[ 0 ] + ".pptx" 

            prs.save( fileLocation ) 
            print( "File saved at location: " + fileLocation )
            file.close()
    else:
        print( "Usage: {0} <file1> <file2> ... <fileN>".format( argv[ 0 ] ) )



def parser( prs, file ):
    block = ""
    for line in file:
        line = line.strip()            # Clean whitespace in line
        if line == "" or line == "\n":          # skip empty lines
            continue

        if line[ 0 ] != '-':
            block += ( "\nSlide Title: " + line )
        else:
            count = 0
            while ( line[ count ] == '-' ): # Get num of initial dashes
                count += 1
            line = line.replace( line[:count], "" ).strip()
            if count == 1:
                print( block )
                block = "\nNew Block:"
                block += ( "\nHeader: " + line )
            elif count == 2:
                block += ( "\n\tTopic: " + line )
            elif count == 3:
                block += ( "\n\t\tBullet: " + line )
            else:
                block += ( "\n\t\t\tSub-Bullet: " + line )
    print( block ) # Catch the last block




if __name__ == "__main__":
    main()
